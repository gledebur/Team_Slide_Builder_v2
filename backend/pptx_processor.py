import os
import logging
from typing import List, Dict, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import tempfile
import shutil

logger = logging.getLogger(__name__)

# Define Consulting Purple color (RGB values)
CONSULTING_PURPLE = RGBColor(102, 45, 145)  # Deep purple color typically used in consulting
GRAY_TEXT = RGBColor(64, 64, 64)  # Dark gray for secondary text
BLACK_TEXT = RGBColor(0, 0, 0)  # Black for primary text

class PowerPointProcessor:
    def __init__(self, cvs_folder: str, output_folder: str, examples_folder: str):
        self.cvs_folder = cvs_folder
        self.output_folder = output_folder
        self.examples_folder = examples_folder
        
    def find_cv_file(self, consultant_name: str) -> Optional[str]:
        """
        Find CV file for a consultant name using flexible matching
        """
        # Get all .pptx files in the CVs folder
        try:
            cv_files = [f for f in os.listdir(self.cvs_folder) if f.endswith('.pptx') and not f.startswith('CV_Placeholder')]
        except OSError:
            logger.error(f"Could not list files in {self.cvs_folder}")
            return None
        
        # First try exact match with standard naming rule
        standard_filename = consultant_name.replace(" ", "_").replace("-", "") + ".pptx"
        if standard_filename in cv_files:
            logger.info(f"Found exact match for {consultant_name}: {standard_filename}")
            return standard_filename
        
        # Try flexible matching - look for files that contain the consultant's name parts
        name_parts = consultant_name.lower().split()
        
        for cv_file in cv_files:
            # Remove .pptx extension and convert to lowercase for comparison
            file_base = cv_file[:-5].lower()  # Remove ".pptx"
            
            # Check if all name parts are present in the filename
            if all(part in file_base for part in name_parts):
                logger.info(f"Found fuzzy match for {consultant_name}: {cv_file}")
                return cv_file
        
        # If no match found, log available files for debugging
        logger.warning(f"No CV file found for {consultant_name}. Available files: {cv_files}")
        return None
        
    def extract_consultant_data_from_template(self, cv_filepath: str, consultant_name: str) -> Dict:
        """
        Extract consultant data from a CV PowerPoint file using the CV_Placeholder structure as reference
        Returns: {
            'first_name': str,
            'last_name': str, 
            'office': str,
            'experience_bullets': List[str],
            'headshot_image': bytes or None
        }
        """
        logger.info(f"Extracting data from {cv_filepath} using template structure")
        
        try:
            prs = Presentation(cv_filepath)
            
            # Assume we're working with the first slide
            if len(prs.slides) == 0:
                raise ValueError(f"No slides found in {cv_filepath}")
                
            slide = prs.slides[0]
            
            # Initialize data
            consultant_data = {
                'first_name': consultant_name.split()[0] if consultant_name else "First",
                'last_name': consultant_name.split()[-1] if consultant_name and len(consultant_name.split()) > 1 else "Last",
                'office': "Global",
                'experience_bullets': [],
                'headshot_image': None
            }
            
            # Find top-left textbox and top-left image
            top_left_text_shape = None
            top_left_image_shape = None
            min_position = float('inf')
            min_image_position = float('inf')
            
            # Extract data from shapes based on CV_Placeholder structure
            for i, shape in enumerate(slide.shapes):
                try:
                    # Find top-left image (headshot)
                    if hasattr(shape, 'image') and consultant_data['headshot_image'] is None:
                        # Calculate position (top + left for simple ranking)
                        position = shape.top + shape.left
                        if position < min_image_position:
                            min_image_position = position
                            top_left_image_shape = shape
                    
                    # Find top-left textbox
                    elif hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        position = shape.top + shape.left
                        
                        # Check if this contains name/position info and is positioned in top-left area
                        if (any(keyword in text.lower() for keyword in ['position', 'office', 'location', 'germany', 'london', 'new york', 'paris', 'berlin', 'zurich', 'geneva', 'munich']) 
                            and position < min_position):
                            min_position = position
                            top_left_text_shape = shape
                        
                        # Also check for "Selected consulting engagement experience" section
                        if 'consulting engagement experience' in text.lower() or 'consulting experience' in text.lower():
                            lines = [line.strip() for line in text.split('\n') if line.strip()]
                            bullets = []
                            
                            for line in lines:
                                # Skip header lines
                                if 'consulting engagement experience' in line.lower() or 'take 3 bullet' in line.lower():
                                    continue
                                    
                                # Look for actual bullet points (meaningful content lines)
                                if len(line) > 20 and not line.startswith('Take '):  # Avoid instruction text
                                    # Clean up bullet formatting
                                    clean_line = line.lstrip('•-▪◦→ ').strip()
                                    if clean_line and len(clean_line) > 20:
                                        bullets.append(clean_line)
                            
                            # Take only first 3 bullets as required
                            consultant_data['experience_bullets'] = bullets[:3]
                            logger.info(f"Extracted {len(consultant_data['experience_bullets'])} experience bullets")
                        
                except Exception as e:
                    logger.warning(f"Error processing shape {i}: {str(e)}")
                    continue
            
            # Extract headshot image
            if top_left_image_shape:
                image_stream = io.BytesIO(top_left_image_shape.image.blob)
                consultant_data['headshot_image'] = image_stream.getvalue()
                logger.info(f"Extracted headshot image from top-left position")
            
            # Extract name and office from top-left textbox
            if top_left_text_shape:
                text = top_left_text_shape.text.strip()
                lines = [line.strip() for line in text.split('\n') if line.strip()]
                
                # Find the name line - usually contains comma and proper name structure
                for line in lines:
                    if ',' in line and any(char.isalpha() for char in line):
                        name_parts = line.split(',')
                        if len(name_parts) >= 2:
                            # Format: "Last Name, First Name" or similar
                            last_name = name_parts[0].strip()
                            first_name = name_parts[1].strip()
                            # Only update if this looks like a proper name (not random text)
                            if len(last_name) < 50 and len(first_name) < 50 and not any(keyword in line.lower() for keyword in ['university', 'msc', 'ba', 'phd', 'degree']):
                                consultant_data['first_name'] = first_name
                                consultant_data['last_name'] = last_name
                                break
                
                # Look for office location in all lines
                for line in lines:
                    if any(keyword in line.lower() for keyword in ['germany', 'london', 'new york', 'paris', 'berlin', 'zurich', 'geneva', 'munich']) and len(line) < 100:
                        consultant_data['office'] = line.strip()
                        break
            
            # Ensure we have exactly 3 bullet points
            while len(consultant_data['experience_bullets']) < 3:
                consultant_data['experience_bullets'].append("Proven track record in client engagement and project delivery")
            
            consultant_data['experience_bullets'] = consultant_data['experience_bullets'][:3]
            
            logger.info(f"Extracted data - First Name: {consultant_data['first_name']}, Last Name: {consultant_data['last_name']}, "
                       f"Office: {consultant_data['office']}, Bullets: {len(consultant_data['experience_bullets'])}")
            
            return consultant_data
            
        except Exception as e:
            logger.error(f"Error extracting data from {cv_filepath}: {str(e)}")
            raise
    
    def _crop_and_resize_image(self, image_bytes: bytes, target_width: int, target_height: int) -> bytes:
        """
        Crop and resize image to fit exactly into the designated space
        """
        try:
            # Load image
            image = Image.open(io.BytesIO(image_bytes))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Calculate crop box for center crop
            img_width, img_height = image.size
            aspect_ratio = target_width / target_height
            img_aspect_ratio = img_width / img_height
            
            if img_aspect_ratio > aspect_ratio:
                # Image is wider than target, crop width
                new_width = int(img_height * aspect_ratio)
                left = (img_width - new_width) // 2
                crop_box = (left, 0, left + new_width, img_height)
            else:
                # Image is taller than target, crop height
                new_height = int(img_width / aspect_ratio)
                top = (img_height - new_height) // 2
                crop_box = (0, top, img_width, top + new_height)
            
            # Crop and resize
            cropped_image = image.crop(crop_box)
            resized_image = cropped_image.resize((target_width, target_height), Image.Resampling.LANCZOS)
            
            # Save to bytes
            output = io.BytesIO()
            resized_image.save(output, format='JPEG', quality=95)
            return output.getvalue()
            
        except Exception as e:
            logger.warning(f"Error processing image: {str(e)}")
            return image_bytes  # Return original if processing fails

    def create_team_slide(self, names: List[str]) -> str:
        """
        Create a team slide using Output_Example_Placeholder_Logic.pptx template.
        Only modifies designated placeholders while preserving all other formatting.
        
        Args:
            names: List of consultant names to include in the team slide (max 4)
            
        Returns:
            Path to the generated Team_Slide_Output.pptx file
        """
        logger.info(f"Creating team slide for consultants: {names}")
        
        # Find and extract data from consultant CV files
        consultants_data = []
        for name in names:
            filename = self.find_cv_file(name)
            if not filename:
                logger.warning(f"CV file not found for {name}, using placeholder data")
                # Create placeholder data for missing consultant
                consultants_data.append({
                    'first_name': name.split()[0] if name else "First",
                    'last_name': name.split()[-1] if name and len(name.split()) > 1 else "Last",
                    'office': "Global",
                    'experience_bullets': [
                        "Extensive experience in strategic consulting",
                        "Proven track record in client engagement",
                        "Specialized in project delivery and transformation"
                    ],
                    'headshot_image': None
                })
            else:
                cv_filepath = os.path.join(self.cvs_folder, filename)
                try:
                    data = self.extract_consultant_data_from_template(cv_filepath, name)
                    consultants_data.append(data)
                except Exception as e:
                    logger.error(f"Failed to extract data from {filename}: {str(e)}")
                    # Use placeholder data if extraction fails
                    consultants_data.append({
                        'first_name': name.split()[0] if name else "First",
                        'last_name': name.split()[-1] if name and len(name.split()) > 1 else "Last",
                        'office': "Global",
                        'experience_bullets': [
                            "Extensive experience in strategic consulting",
                            "Proven track record in client engagement", 
                            "Specialized in project delivery and transformation"
                        ],
                        'headshot_image': None
                    })
        
        # Load the output template (fix the typo in filename)
        template_path = os.path.join(self.examples_folder, 'Output_Example_Placeholder_Logic.pptx')
        
        # Check for the typo version first (current file), then the correct name
        if not os.path.exists(template_path):
            template_path = os.path.join(self.examples_folder, 'Outpout_Example_Placeholder_Logic.pptx')
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Output template not found in {self.examples_folder}")
        
        logger.info(f"Loading output template from {template_path}")
        prs = Presentation(template_path)
        slide = prs.slides[0]
        
        # First, collect all placeholder shapes organized by position/proximity
        text_placeholders = []
        image_placeholders = []
        experience_shapes = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape_text = shape.text.strip()
                
                # Find text placeholders
                if shape_text in ["First Name", "Last Name", "Office"]:
                    text_placeholders.append({
                        'shape': shape,
                        'placeholder_type': shape_text,
                        'position': shape.top + shape.left  # Simple position ranking
                    })
                
                # Find experience shapes
                elif "years of consulting experience" in shape_text.lower():
                    experience_shapes.append({
                        'shape': shape,
                        'position': shape.top + shape.left
                    })
            
            # Find image placeholders
            should_check_image = False
            if hasattr(shape, 'text') and "replace picture" in shape.text.lower():
                should_check_image = True
            elif hasattr(shape, 'image') and hasattr(shape, 'name'):
                if "replace" in shape.name.lower() or "picture" in shape.name.lower():
                    should_check_image = True
            
            if should_check_image:
                image_placeholders.append({
                    'shape': shape,
                    'position': shape.top + shape.left
                })
        
        # Sort placeholders by position (top-left to bottom-right order)
        text_placeholders.sort(key=lambda x: x['position'])
        image_placeholders.sort(key=lambda x: x['position'])
        experience_shapes.sort(key=lambda x: x['position'])
        
        # Group text placeholders by consultant (every 3 placeholders = 1 consultant)
        consultant_text_groups = []
        for i in range(0, len(text_placeholders), 3):
            group = text_placeholders[i:i+3]
            if len(group) == 3:  # Only process complete groups
                consultant_text_groups.append(group)
        
        # Process up to 4 consultants
        shapes_to_remove = []
        new_images = []
        
        for i, consultant_data in enumerate(consultants_data[:4]):
            try:
                # Update text placeholders for this consultant
                if i < len(consultant_text_groups):
                    for placeholder in consultant_text_groups[i]:
                        shape = placeholder['shape']
                        placeholder_type = placeholder['placeholder_type']
                        
                        if placeholder_type == "First Name":
                            shape.text = consultant_data['first_name']
                            logger.info(f"Replaced 'First Name' with '{consultant_data['first_name']}' for consultant {i+1}")
                        elif placeholder_type == "Last Name":
                            shape.text = consultant_data['last_name']
                            logger.info(f"Replaced 'Last Name' with '{consultant_data['last_name']}' for consultant {i+1}")
                        elif placeholder_type == "Office":
                            shape.text = consultant_data['office']
                            logger.info(f"Replaced 'Office' with '{consultant_data['office']}' for consultant {i+1}")
                
                # Add bullets to experience shape for this consultant
                if i < len(experience_shapes):
                    experience_shape = experience_shapes[i]['shape']
                    existing_text = experience_shape.text.rstrip()
                    bullets_text = "\n".join([f"• {bullet}" for bullet in consultant_data['experience_bullets']])
                    experience_shape.text = existing_text + "\n\n" + bullets_text
                    logger.info(f"Added experience bullets for consultant {i+1}")
                
                # Replace image placeholder for this consultant
                if i < len(image_placeholders) and consultant_data['headshot_image']:
                    placeholder_info = image_placeholders[i]
                    shape = placeholder_info['shape']
                    
                    # Get position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    # Process the headshot image
                    target_width = int(width.inches * 96)  # Convert to pixels (96 DPI)
                    target_height = int(height.inches * 96)
                    
                    processed_image = self._crop_and_resize_image(
                        consultant_data['headshot_image'],
                        target_width,
                        target_height
                    )
                    
                    # Save processed image temporarily
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                        temp_img.write(processed_image)
                        temp_img_path = temp_img.name
                    
                    # Mark for removal and schedule new image
                    shapes_to_remove.append(shape)
                    new_images.append({
                        'path': temp_img_path,
                        'left': left,
                        'top': top,
                        'width': width,
                        'height': height
                    })
                    
                    logger.info(f"Scheduled image replacement for consultant {i+1}")
                
            except Exception as e:
                logger.error(f"Failed to update consultant {i+1} data: {str(e)}")
                continue
        
        # Remove old placeholder image shapes
        for shape in shapes_to_remove:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                logger.warning(f"Failed to remove placeholder shape: {str(e)}")
        
        # Add new images
        for img_info in new_images:
            try:
                slide.shapes.add_picture(
                    img_info['path'],
                    img_info['left'],
                    img_info['top'],
                    img_info['width'],
                    img_info['height']
                )
                # Clean up temp file
                os.unlink(img_info['path'])
            except Exception as e:
                logger.warning(f"Failed to add new image: {str(e)}")
                # Clean up temp file even on failure
                try:
                    os.unlink(img_info['path'])
                except:
                    pass
        
        # Save the final presentation
        output_filename = "Team_Slide_Output.pptx"
        output_path = os.path.join(self.output_folder, output_filename)
        prs.save(output_path)
        
        logger.info(f"Team slide saved to {output_path}")
        return output_path

    # Keep the old method name for backward compatibility
    def generate_team_slide(self, consultant_names: List[str], filenames: List[str] = None) -> str:
        """
        Backward compatibility wrapper for create_team_slide
        """
        return self.create_team_slide(consultant_names)